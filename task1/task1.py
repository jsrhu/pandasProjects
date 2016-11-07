# -*- coding: utf-8 -*-
"""
Created on Mon Sep 26 12:51:14 2016

@author: jsrhu

STATUS:
TODO
- change path to output to ignored directory
"""
import sys
import os
import pandas as pd

import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches

from urlparse import urlparse

compressed_data = '../data/csv/backtest_alpha_for_ernest.csv.gz'

headers = ['story_sentiment','story_volume','story_traffic','story_shares','article_sentiment','article_traffic','event_impact_score_overall','event_impact_score_entity_1','event_impact_score_entity_2','avg_day_sentiment']

plot_type = 'max_min'

'''
Reads gzip compressed csv file and converts date string into date-time object
Returns dataframerepresentation of csv file
'''
def readLimitedRows(csv_file, rows):
    date_format="%Y-%m-%d %H:%M:%S %Z"
    dateparse = lambda x: pd.datetime.strptime(x, date_format)
    return pd.read_csv(filepath_or_buffer=csv_file, compression='gzip', nrows=rows, header=0, date_parser=dateparse, parse_dates=["harvested_at"])

'''
Reads gzip compressed csv file and converts date string into date-time object
Does not return anything
'''
def readFull(csv_file):
    date_format="%Y-%m-%d %H:%M:%S %Z"
    dateparse = lambda x: pd.datetime.strptime(x, date_format)
    return pd.read_csv(filepath_or_buffer=csv_file, compression='gzip', header=0, date_parser=dateparse, parse_dates=["harvested_at"])

'''
Returns dataframe with select specififed columns from the original
'''
def selectColumnsDataFrame(data_frame,columns):
    data_frame = data_frame[columns]
    return data_frame

'''
increases the y scale of a pyplot instance by a given decimal percentage
'''
def increaseYScaleRange(plot, dataframe_min, dataframe_max, column, percentage):
    plot.ylim(dataframe_min[column].min()-abs(percentage*dataframe_min[column].min()),dataframe_max[column].max()+abs(percentage*dataframe_max[column].max()))

'''
Attempts to create directory for plot images
Returns path to directory
'''
def createDir(name, descriptor):
    path = name+'_'+descriptor
    try:
        os.mkdir(path)
        return path
    except:
        return path

def addImageSlide(presentation, layout, image_path):
    slide_layout = presentation.slide_layouts[layout]

    slide = presentation.slides.add_slide(slide_layout)

    left = Inches(1)
    top = Inches(1)
    height = Inches(5.5)
    return slide.shapes.add_picture(image_file=image_path, left=left, top=top, height=height)

def savePresentation(presentation, title):
    presentation.save(title+'.pptx')

def myurlparse(x):    
    try:
        return urlparse(x).netloc
    except:
        return x
    
def main():
    rows=1000000
    df = readLimitedRows(csv_file=compressed_data,rows=rows)
    #df = readFull(csv_file=compressed_data)
    last_date = df['harvested_at'].max().date()
    df_max = pd.DataFrame()
    df_min = pd.DataFrame()
    df_group = df.groupby( df['harvested_at'].apply(lambda x: x.date() ) )
    try:
        output = os.mkdir('output')
    except:
        output = 'output'
    os.chdir(output)

    '''
    # TICKER COUNT
    ticker_count_1 = df_group['entities_ticker_1'].nunique()
    ticker_count_1.plot(title='Number of Primary Tickers from Harvested Articles by Day', rot=45, figsize=(11,8), x_compat=True)
    plt.tight_layout()
    plt.savefig('unique_ticker1_plot_'+str(rows)+'rows.png')
    plt.clf()

    ticker_count_2 = df_group['entities_ticker_2'].nunique()
    ticker_count_2.plot(title='Number of Secondary Tickers from Harvested Articles by Day', rot=45, figsize=(11,8), x_compat=True)
    plt.tight_layout()
    plt.savefig('unique_ticker2_plot_'+str(rows)+'rows.png')
    plt.clf()
    '''
    
    '''
    # MAXMIN PLOTS
    #should look for a better way to do this
    #perhaps new df called df_minmax_$group where $group is one of the numeric columns
    for key,group in df_group:
        group_max = group.max(numeric_only=True)
        group_min = group.min(numeric_only=True)
        group_max['date'] = key
        group_min['date'] = key

        df_max = pd.concat([df_max, group_max.to_frame().T])
        df_min = pd.concat([df_min, group_min.to_frame().T])

    df_max.set_index(keys='date', inplace=True)
    df_min.set_index(keys='date', inplace=True)

    df_max = selectColumnsDataFrame(data_frame=df_max,columns=headers)
    df_min = selectColumnsDataFrame(data_frame=df_min,columns=headers)

    plot_dir = createDir(name=plot_type, descriptor='plots')
    os.chdir(plot_dir)

    presentation = Presentation()
    descriptor = 'max_min_plot'
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[layout_title])
    title_slide.shapes.title.text = 'Maximum and Minimum Plots of Quantitative Metrics in:\n'+'Accern Data'

    for column in df_max:
        column_plot = df_max[column].plot()
        column_plot = df_min[column].plot(ax=column_plot, title=column, x_compat=True, rot=45, figsize=(11,8))

        increaseYScaleRange(plt, dataframe_min=df_min, dataframe_max=df_max, column=column, percentage=.10)

        image_path = column+'_'+descriptor+'.png'
        plt.tight_layout()
        plt.savefig(image_path)
        plt.clf()

        addImageSlide(presentation=presentation, layout=layout_blank, image_path=image_path)

    os.chdir('..')
    ppt_dir = createDir(name='ppt', descriptor=descriptor)
    os.chdir(ppt_dir)
    savePresentation(presentation,descriptor)
    '''
    return 0

if __name__ == "__main__":
    sys.exit(main())
