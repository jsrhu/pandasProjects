# -*- coding: utf-8 -*-
"""
Created on Fri Sep 30 10:15:42 2016

@author: jsrhu
"""

import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

layout_title = 0
layout_title_content = 1
layout_section_header = 2
layout_two_content = 3
layout_comparison = 4
layout_title_only = 5
layout_blank = 6
layout_content_caption = 7
layout_picture_caption = 8
    
def addSlide():
    img_path = 'test.png'
    
    prs = Presentation()
    
    slide_layout = prs.slide_layouts[layout_picture_caption]
    
    plt.savefig(img_path)
    plt.show()
    
    slide = prs.slides.add_slide(slide_layout)
    
    left = Inches(.8)
    top = Inches(.8)
    height = Inches(5.5)
    pic = slide.shapes.add_picture(image_file=img_path, left=left, top=top, height=height)
    
    prs.save('test.pptx')